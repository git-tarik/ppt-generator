import logging
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE
from pptx.util import Inches
import io

logger = logging.getLogger("SlideBuilder")

def update_slide_content(slide, slide_data: dict, template_images: dict = None):
    """
    Updates the text content of a slide's placeholders and adds images from template.
    
    Args:
        slide: The slide object to update
        slide_data: Dictionary containing title, bullets, and notes
        template_images: Dictionary containing categorized images from template
    """
    
    # 1. Identify Placeholders
    title_ph = None
    body_ph = None
    
    # Scan shapes to find placeholders (even if cloned)
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
            
        ph_type = shape.placeholder_format.type
        
        # TITLE
        if ph_type == PP_PLACEHOLDER.TITLE or ph_type == PP_PLACEHOLDER.CENTER_TITLE:
            if not title_ph: title_ph = shape
            
        # BODY
        if ph_type == PP_PLACEHOLDER.BODY or ph_type == PP_PLACEHOLDER.OBJECT:
            if not body_ph: body_ph = shape
            
    # 2. Update Title
    if title_ph and 'title' in slide_data:
        try:
            if title_ph.text_frame:
                if title_ph.text_frame.paragraphs:
                    # Preserve formatting of the first paragraph if it exists
                    title_ph.text_frame.paragraphs[0].text = slide_data['title']
                else:
                    title_ph.text = slide_data['title']
        except Exception as e:
            logger.warning(f"Failed to update title: {e}")

    # 3. Update Body (Bullets)
    if body_ph and 'bullets' in slide_data:
        try:
            text_frame = body_ph.text_frame
            text_frame.clear() 
            
            bullets = slide_data.get("bullets", [])
            for bullet_text in bullets:
                p = text_frame.add_paragraph()
                p.text = bullet_text
                p.level = 0
        except Exception as e:
            logger.warning(f"Failed to update body: {e}")
            
    # 4. Update Notes
    if "notes" in slide_data and slide_data["notes"]:
        try:
            if not slide.has_notes_slide:
                slide.notes_slide
            
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_data["notes"]
        except Exception as e:
            logger.warning(f"Failed to add notes: {e}")
    
    # 5. Add Images from Template (if available)
    if template_images and template_images.get("categorized"):
        try:
            add_template_images_to_slide(slide, template_images)
        except Exception as e:
            logger.warning(f"Failed to add template images: {e}")


def add_template_images_to_slide(slide, template_images: dict):
    """
    Adds images from the template to the slide.
    Prioritizes logos and reuses them in consistent positions.
    """
    categorized = template_images.get("categorized", {})
    
    # Add logos (typically small images in corners)
    logos = categorized.get("logos", [])
    for logo in logos[:2]:  # Limit to 2 logos to avoid clutter
        try:
            add_image_to_slide(slide, logo)
            logger.debug(f"Added logo to slide")
        except Exception as e:
            logger.debug(f"Could not add logo: {e}")
    
    # Optionally add one background or content image if space allows
    # This is conservative to avoid overcrowding slides
    backgrounds = categorized.get("backgrounds", [])
    if backgrounds and len(logos) == 0:  # Only if no logos were added
        try:
            add_image_to_slide(slide, backgrounds[0])
            logger.debug(f"Added background image to slide")
        except Exception as e:
            logger.debug(f"Could not add background: {e}")


def add_image_to_slide(slide, image_data: dict):
    """
    Adds a single image to a slide at its original position.
    """
    try:
        blob = image_data.get("blob")
        position = image_data.get("position", {})
        
        if not blob:
            return
        
        # Create image stream
        image_stream = io.BytesIO(blob)
        
        # Add picture to slide at original position
        left = position.get("left", Inches(1))
        top = position.get("top", Inches(1))
        width = position.get("width", Inches(1))
        height = position.get("height", Inches(1))
        
        slide.shapes.add_picture(image_stream, left, top, width, height)
        
    except Exception as e:
        logger.debug(f"Failed to add individual image: {e}")
        raise
