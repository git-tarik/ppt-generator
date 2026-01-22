from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io
import logging
from typing import List, Dict, Any
from PIL import Image

logger = logging.getLogger("ImageExtractor")
logger.setLevel(logging.INFO)

def extract_images_from_template(template_bytes: bytes) -> Dict[str, Any]:
    """
    Extracts all images from a PowerPoint template.
    Returns a dictionary containing image data and metadata.
    """
    try:
        prs = Presentation(io.BytesIO(template_bytes))
    except Exception as e:
        logger.error(f"Failed to load template for image extraction: {e}")
        return {"images": [], "error": str(e)}
    
    images_catalog = []
    image_id = 0
    
    # Iterate through all slides
    for slide_idx, slide in enumerate(prs.slides):
        # Check all shapes on the slide
        for shape_idx, shape in enumerate(slide.shapes):
            try:
                # Check if shape is a picture
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_data = extract_image_from_shape(shape, slide_idx, shape_idx, image_id)
                    if image_data:
                        images_catalog.append(image_data)
                        image_id += 1
                        
                # Check if shape is a placeholder that might contain an image
                elif hasattr(shape, 'image'):
                    image_data = extract_image_from_shape(shape, slide_idx, shape_idx, image_id)
                    if image_data:
                        images_catalog.append(image_data)
                        image_id += 1
                        
            except Exception as e:
                logger.debug(f"Could not extract image from shape {shape_idx} on slide {slide_idx}: {e}")
                continue
    
    logger.info(f"Extracted {len(images_catalog)} images from template")
    
    return {
        "images": images_catalog,
        "total_count": len(images_catalog)
    }


def extract_image_from_shape(shape, slide_idx: int, shape_idx: int, image_id: int) -> Dict[str, Any]:
    """
    Extracts image data from a shape.
    Returns metadata including position, size, and binary data.
    """
    try:
        # Get image binary data
        image_blob = shape.image.blob
        
        # Get image format
        image_format = shape.image.content_type.split('/')[-1]
        
        # Get position and size
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
        
        # Try to get image dimensions using Pillow
        try:
            img = Image.open(io.BytesIO(image_blob))
            img_width, img_height = img.size
        except:
            img_width, img_height = None, None
        
        return {
            "id": image_id,
            "slide_index": slide_idx,
            "shape_index": shape_idx,
            "blob": image_blob,
            "format": image_format,
            "position": {
                "left": left,
                "top": top,
                "width": width,
                "height": height
            },
            "original_dimensions": {
                "width": img_width,
                "height": img_height
            }
        }
        
    except Exception as e:
        logger.debug(f"Failed to extract image metadata: {e}")
        return None


def categorize_images(images_catalog: List[Dict[str, Any]]) -> Dict[str, List[Dict[str, Any]]]:
    """
    Categorizes images into types: logos, backgrounds, content images.
    This is a heuristic-based approach.
    """
    logos = []
    backgrounds = []
    content_images = []
    
    for img in images_catalog:
        pos = img.get("position", {})
        width = pos.get("width", 0)
        height = pos.get("height", 0)
        top = pos.get("top", 0)
        left = pos.get("left", 0)
        
        # Heuristic: Small images in corners are likely logos
        # PowerPoint uses EMUs (914400 EMUs = 1 inch)
        # Typical slide is ~10 inches wide, ~7.5 inches tall
        slide_width = 9144000  # ~10 inches in EMUs
        slide_height = 6858000  # ~7.5 inches in EMUs
        
        # Logo heuristic: small image (< 15% of slide area) in top corners
        area_ratio = (width * height) / (slide_width * slide_height) if slide_width and slide_height else 0
        is_small = area_ratio < 0.15
        is_top = top < slide_height * 0.2
        is_corner = left < slide_width * 0.2 or left > slide_width * 0.8
        
        if is_small and is_top and is_corner:
            logos.append(img)
        # Background heuristic: large image covering most of slide
        elif area_ratio > 0.7:
            backgrounds.append(img)
        else:
            content_images.append(img)
    
    logger.info(f"Categorized images: {len(logos)} logos, {len(backgrounds)} backgrounds, {len(content_images)} content")
    
    return {
        "logos": logos,
        "backgrounds": backgrounds,
        "content": content_images
    }
