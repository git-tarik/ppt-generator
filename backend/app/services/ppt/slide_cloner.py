from pptx import Presentation
import copy
import logging

logger = logging.getLogger("SlideCloner")

def clone_slide(pres, source_slide):
    """
    Duplicate a slide in a presentation by deep-copying its XML.
    Adds the new slide to the end of the presentation.
    """
    try:
        # 1. Create a new slide based on the same layout
        # This gives us a container linked to the correct layout
        new_slide = pres.slides.add_slide(source_slide.slide_layout)
        
        # 2. Clear default shapes
        # We assume the source slide has the 'true' state of the placeholders/shapes we want.
        # We must clear the default ones created by add_slide so we don't get duplicates.
        for shape in list(new_slide.shapes):
             sp = shape.element
             sp.getparent().remove(sp)
             
        # 3. Copy shapes from source to new
        # Append deeper copies of source shapes to the new slide
        for shape in source_slide.shapes:
            new_el = copy.deepcopy(shape.element)
            new_slide.shapes._spTree.append(new_el)
            
        return new_slide

    except Exception as e:
        logger.error(f"Slide cloning failed: {e}")
        # Fail safe: return a fresh slide (may lose content but prevents crash)
        return pres.slides.add_slide(source_slide.slide_layout)
