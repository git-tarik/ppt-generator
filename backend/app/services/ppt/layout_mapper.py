from pptx.presentation import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
import logging

logger = logging.getLogger("LayoutMapper")

def map_layout(prs: Presentation, slide_type: str) -> int:
    """
    Selects the best layout index based on slide intent.
    Scans layouts for specific placeholders.
    """
    best_layout_idx = 0
    
    # Define what we are looking for
    # TITLE_SLIDE: Title (1) + optionally Subtitle (4)
    # CONTENT_SLIDE: Title (1) + Body (2, 14, etc - usually type 2 or 7 for object)
    
    # Iterate through all layouts to find the best match
    for i, layout in enumerate(prs.slide_layouts):
        has_title = False
        has_body = False
        
        for shape in layout.placeholders:
            ph_type = shape.placeholder_format.type
            
            if ph_type == PP_PLACEHOLDER.TITLE or ph_type == PP_PLACEHOLDER.CENTER_TITLE:
                has_title = True
            
            if ph_type == PP_PLACEHOLDER.BODY or ph_type == PP_PLACEHOLDER.OBJECT:
                has_body = True
        
        if slide_type == "TITLE_SLIDE":
            # Prefer Title-only or Title+Subtitle. try to avoid Body.
            if has_title and not has_body:
                return i
            # If we find a title layout (often index 0), keep it as candidate
            if has_title and i == 0:
                best_layout_idx = i

        elif slide_type == "CONTENT_SLIDE":
            # Must have Title AND Body
            if has_title and has_body:
                return i
    
    # Fallback Logic if no perfect match found
    if slide_type == "CONTENT_SLIDE":
        # If we didn't find a layout with body, try to find one that at least looks like a content slide (usually index 1)
        if len(prs.slide_layouts) > 1:
            return 1
            
    return best_layout_idx

def find_placeholders(slide_layout):
    """
    Helper to identify placeholders in a layout.
    Returns dict like { 'title': shape, 'body': shape }
    """
    placeholders = {}
    for shape in slide_layout.placeholders:
        ph_type = shape.placeholder_format.type
        
        # PP_PLACEHOLDER.TITLE = 1, CENTER_TITLE = 3
        if ph_type == PP_PLACEHOLDER.TITLE or ph_type == PP_PLACEHOLDER.CENTER_TITLE:
            placeholders['title'] = shape
        
        # PP_PLACEHOLDER.BODY = 2, OBJECT = 7
        # We assign the first body found to 'body'
        if (ph_type == PP_PLACEHOLDER.BODY or ph_type == PP_PLACEHOLDER.OBJECT) and 'body' not in placeholders:
            placeholders['body'] = shape
            
    return placeholders
