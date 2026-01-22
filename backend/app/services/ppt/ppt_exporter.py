from pptx import Presentation
import io
import logging
from .slide_cloner import clone_slide
from .slide_builder import update_slide_content

logger = logging.getLogger("PPTExporter")
logger.setLevel(logging.INFO)

def generate_presentation(template_content: bytes, slide_plan: dict, template_metadata: dict = None) -> io.BytesIO:
    """
    Generates a PPTX file by cloning slides from the template.
    Strictly follows slide plan count and reuses template images.
    
    Args:
        template_content: Binary content of the template PPTX
        slide_plan: Dictionary containing slides and metadata
        template_metadata: Optional metadata including images, colors, fonts
    """
    try:
        prs = Presentation(io.BytesIO(template_content))
    except Exception as e:
        logger.error(f"Failed to load template: {e}")
        raise ValueError("Invalid template file")

    slides_data = slide_plan.get("slides", [])
    if not slides_data:
        raise ValueError("Empty slide plan")

    # Capture original templates
    template_slides = list(prs.slides)
    if not template_slides:
        raise ValueError("Template has no slides to clone from")

    num_template_slides = len(template_slides)
    
    # Extract template images if available
    template_images = None
    if template_metadata and template_metadata.get("images"):
        template_images = template_metadata["images"]
        logger.info(f"Using {template_images.get('total', 0)} images from template")
    
    # 3. Generate New Slides
    # Loop ONLY over slide_plan["slides"]
    for i, slide_data in enumerate(slides_data):
        logger.debug(f"Creating slide {i+1}/{len(slides_data)}")
        
        # Strict modulo mapping
        base_slide = template_slides[i % num_template_slides]
        
        # Clone (Must return a NEW object)
        new_slide = clone_slide(prs, base_slide)
        
        # Update with content and images
        update_slide_content(new_slide, slide_data, template_images)

    # 4. Cleanup: Remove the original template slides
    # Iterate backwards through the original count and remove element
    for i in range(num_template_slides - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

    # 5. Export
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    
    logger.info(f"Generated presentation with {len(slides_data)} slides")
    
    return output
