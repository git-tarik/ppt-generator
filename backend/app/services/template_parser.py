from pptx import Presentation
import io
import logging
from typing import Dict, Any
from app.services.ppt.image_extractor import extract_images_from_template, categorize_images

# Configure logger
logger = logging.getLogger("TemplateParser")
logger.setLevel(logging.INFO)

if not logger.handlers:
    handler = logging.StreamHandler()
    formatter = logging.Formatter(
        "[%(name)s] %(levelname)s: %(message)s"
    )
    handler.setFormatter(formatter)
    logger.addHandler(handler)

logger.propagate = False


def analyze_presentation(file_content: bytes) -> dict:
    """
    Analyzes a PPTX file and extracts comprehensive metadata including:
    - Layout information
    - Theme colors and fonts
    - Images catalog
    Returns a dictionary with all extracted information.
    """
    try:
        prs = Presentation(io.BytesIO(file_content))
    except Exception as e:
        logger.error(f"Failed to load presentation: {e}")
        return {"error": "Invalid PPTX file"}

    metadata = {
        "layout_count": len(prs.slide_layouts),
        "layouts": [],
        "theme": {
            "colors": {},
            "fonts": {}
        },
        "image_placeholders": [],
        "images": []
    }

    # Extract layout information
    for index, layout in enumerate(prs.slide_layouts):
        layout_data = {
            "index": index,
            "name": layout.name,
            "placeholders": []
        }
        
        for shape in layout.placeholders:
            ph_type = str(shape.placeholder_format.type)
            layout_data["placeholders"].append(ph_type)
            
            if 'PICTURE' in ph_type or 'BITMAP' in ph_type: 
                metadata["image_placeholders"].append({
                    "layout_index": index,
                    "placeholder_type": ph_type
                })

        metadata["layouts"].append(layout_data)

    # Extract theme colors
    try:
        theme_colors = extract_theme_colors(prs)
        metadata["theme"]["colors"] = theme_colors
    except Exception as e:
        logger.warning(f"Could not extract theme colors: {e}")
        metadata["theme"]["colors"] = {"note": "Color extraction failed"}

    # Extract theme fonts
    try:
        theme_fonts = extract_theme_fonts(prs)
        metadata["theme"]["fonts"] = theme_fonts
    except Exception as e:
        logger.warning(f"Could not extract theme fonts: {e}")
        metadata["theme"]["fonts"] = {"note": "Font extraction failed"}

    # Extract images from template
    try:
        images_data = extract_images_from_template(file_content)
        if images_data.get("images"):
            categorized = categorize_images(images_data["images"])
            metadata["images"] = {
                "total": images_data.get("total_count", 0),
                "categorized": categorized,
                "raw": images_data["images"]
            }
    except Exception as e:
        logger.warning(f"Could not extract images: {e}")
        metadata["images"] = {"total": 0, "error": str(e)}

    logger.info(f"Detected {metadata['layout_count']} layouts")
    logger.info(f"Extracted {metadata['images'].get('total', 0)} images")
    
    for layout in metadata["layouts"]:
        logger.info(f"Layout {layout['index']} '{layout['name']}': {layout['placeholders']}")
        
    return metadata


def extract_theme_colors(prs: Presentation) -> Dict[str, Any]:
    """
    Extracts theme colors from the presentation.
    Returns a dictionary of color names and their RGB values.
    """
    colors = {}
    
    try:
        # Access the theme part
        theme = prs.part.package.part_related_by(
            'http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme'
        )
        
        # Parse theme XML for color scheme
        theme_xml = theme.blob
        # Note: Full XML parsing would be complex, so we'll extract basic info
        
        # For now, extract colors from actual slides as a fallback
        color_samples = []
        for slide in prs.slides:
            for shape in slide.shapes:
                try:
                    if hasattr(shape, 'fill') and shape.fill.type == 1:  # SOLID fill
                        if hasattr(shape.fill, 'fore_color'):
                            rgb = shape.fill.fore_color.rgb
                            if rgb:
                                color_samples.append({
                                    "rgb": f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                                })
                except:
                    continue
        
        # Get unique colors (limit to first 10)
        unique_colors = []
        seen = set()
        for color in color_samples:
            rgb = color["rgb"]
            if rgb not in seen and len(unique_colors) < 10:
                unique_colors.append(color)
                seen.add(rgb)
        
        colors["extracted_colors"] = unique_colors
        colors["count"] = len(unique_colors)
        
    except Exception as e:
        logger.debug(f"Theme color extraction error: {e}")
        colors["note"] = "Using fallback color extraction"
    
    return colors


def extract_theme_fonts(prs: Presentation) -> Dict[str, Any]:
    """
    Extracts theme fonts from the presentation.
    Returns a dictionary of font information.
    """
    fonts = {}
    
    try:
        # Extract fonts from actual text in slides
        font_samples = set()
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.name:
                                font_samples.add(run.font.name)
        
        fonts["detected_fonts"] = list(font_samples)[:5]  # Limit to 5 most common
        fonts["count"] = len(font_samples)
        
    except Exception as e:
        logger.debug(f"Font extraction error: {e}")
        fonts["note"] = "Font extraction failed"
    
    return fonts

