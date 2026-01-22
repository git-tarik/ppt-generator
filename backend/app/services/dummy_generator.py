from pptx import Presentation
import io

def generate_dummy_pptx(title_text: str, guidance: str | None) -> io.BytesIO:
    prs = Presentation()

    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = title_text or "Demo Slide"
    subtitle.text = f"Guidance used: {guidance}" if guidance else "No guidance provided"

    bullet_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    
    shapes = slide2.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Dummy Content"
    
    tf = body_shape.text_frame
    tf.text = "This is hardcoded text"
    
    p = tf.add_paragraph()
    p.text = "Bullet point 1"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Bullet point 2"
    p.level = 1

    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    
    return pptx_io
